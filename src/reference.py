"""Leitura do PPT de referência do framework.

O documento `reference/Framework_MMM_x_MTA_Digital_Analytics.pptx` é a fonte da
linguagem e da ordem das camadas deste app. Aqui ele é lido com `python-pptx`
para que a Home possa exibir o conteúdo original sem duplicar texto no código.

Tudo é opcional: se o arquivo ou a lib não estiverem presentes, as funções
devolvem vazio e a UI simplesmente não mostra o bloco.
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
import streamlit as st

REFERENCE_DIR = Path(__file__).resolve().parent.parent / "reference"


def reference_path() -> Path | None:
    """Primeiro .pptx encontrado em reference/ (None se não houver)."""
    if not REFERENCE_DIR.exists():
        return None
    files = sorted(REFERENCE_DIR.glob("*.pptx"))
    return files[0] if files else None


@st.cache_data(show_spinner=False)
def load_slides() -> pd.DataFrame:
    """Texto de cada slide do PPT: slide | titulo | secao | texto.

    Heurística de parsing: o primeiro parágrafo curto em caixa alta costuma ser
    a seção (ex.: "03 · CONCEITO") e o seguinte é o título do slide.
    """
    path = reference_path()
    if path is None:
        return pd.DataFrame(columns=["slide", "secao", "titulo", "texto"])

    try:
        from pptx import Presentation
    except Exception:
        return pd.DataFrame(columns=["slide", "secao", "titulo", "texto"])

    try:
        presentation = Presentation(str(path))
    except Exception:
        return pd.DataFrame(columns=["slide", "secao", "titulo", "texto"])

    rows = []
    for i, slide in enumerate(presentation.slides, start=1):
        blocks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    blocks.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        blocks.append(" | ".join(cells))

        if not blocks:
            continue

        secao = ""
        titulo = ""
        for block in blocks[:3]:
            line = block.splitlines()[0].strip()
            if not secao and ("·" in line and line.upper() == line and len(line) < 60):
                secao = line
            elif not titulo and 3 < len(line) < 80:
                titulo = line

        rows.append(
            {
                "slide": i,
                "secao": secao,
                "titulo": titulo or f"Slide {i}",
                "texto": "\n".join(blocks),
            }
        )

    return pd.DataFrame(rows)


@st.cache_data(show_spinner=False)
def reference_available() -> bool:
    return not load_slides().empty


def search_slides(query: str) -> pd.DataFrame:
    """Busca textual simples nos slides (case-insensitive)."""
    slides = load_slides()
    if slides.empty or not query.strip():
        return slides
    mask = slides["texto"].str.contains(query.strip(), case=False, na=False, regex=False)
    return slides[mask]


# ---------------------------------------------------------------------------
# Matriz de decisão (slide 07 · MATRIZ DE DECISÃO)
# ---------------------------------------------------------------------------
DECISION_MATRIX = pd.DataFrame(
    [
        ("Definir orçamento anual/trimestral entre canais", "MMM (Bayesiano, log-log)",
         "Precisa comparar mídias on/offline e captar efeito de longo prazo",
         "💰 Otimizador de Budget"),
        ("Otimizar lances/orçamento diário dentro do Google Ads", "MTA — Data-Driven",
         "Decisão tática de curtíssimo prazo, dado granular disponível",
         "🔀 Modelos de Atribuição"),
        ("Avaliar criativo/público dentro do Meta Ads", "MTA — Data-Driven / Linear",
         "Comparação interna ao canal, mesma fonte de tracking",
         "🎯 Propensão à Conversão"),
        ("Medir contribuição de campanhas de TV/OOH/rádio", "MMM",
         "Não existe tracking individual possível nesses canais",
         "🧪 MMM Modelagem"),
        ("Detectar canibalização entre Branded Search e Social", "MMM (halo/cannibalization)",
         "MTA credita erroneamente o último clique (branded search)",
         "🕸️ Markov e Shapley"),
        ("Jornada de compra longa (B2B, alto ticket)", "MMM + MTA de posição (U-shaped)",
         "Precisa equilibrar descoberta (topo) e decisão (fundo)",
         "🔀 Modelos de Atribuição"),
        ("E-commerce com jornada 100% digital e rastreável", "MTA como principal, MMM como validação",
         "Alta cobertura de tracking permite MTA mais confiável",
         "🕸️ Markov e Shapley"),
        ("Simular 'e se aumentar 20% em TikTok Ads'", "MMM (curvas de saturação)",
         "Cenário de simulação e resposta não-linear ao investimento",
         "💰 Otimizador de Budget"),
        ("Reportar performance de campanha para o time de mídia", "MTA — Data-Driven + Last-Click",
         "Necessário para operação tática do dia a dia",
         "🔀 Modelos de Atribuição"),
        ("Auditoria de budget anual para o CFO/Board", "MMM (contribution + due-to)",
         "Visão holística, defensável estatisticamente, cross-canal",
         "🧪 MMM Modelagem"),
        ("Validar causalmente a incrementalidade de um canal", "Geo-Experiment / Holdout (DiD)",
         "Regressão é correlação; só o experimento prova o incremental",
         "🌍 Geo-Holdout × MMM"),
        ("Decidir se para o teste antes do prazo planejado", "Teste sequencial (SPRT)",
         "Permite parada antecipada sem inflar o falso-positivo",
         "📈 Simulador e Resultados"),
    ],
    columns=["Cenário de negócio", "Técnica recomendada", "Por quê", "Página do app"],
)

# ---------------------------------------------------------------------------
# Números do sumário executivo (slide 01 · EXECUTIVO)
# ---------------------------------------------------------------------------
EXECUTIVE_STATS = [
    ("40–60%", "da jornada é tudo que o MTA tradicional ainda enxerga",
     "era 80–90% antes de iOS 14.5+, cookieless e consent mode"),
    ("14–38%", "de ganho potencial de ROI em mídia",
     "aplicando MMM + MTA de forma combinada (Accenture, 2021)"),
    ("85–95%", "de acurácia possível de um MMM moderno",
     "contra 40–60% de um MTA puro no cenário atual de privacidade"),
]

# ---------------------------------------------------------------------------
# As 7 etapas do Roadmap de Testes A/B (slide 20.2)
# ---------------------------------------------------------------------------
ROADMAP_STAGES = [
    ("Opportunity", "Backlog de oportunidades",
     "MMM aponta driver com contribuição abaixo do esperado · MTA aponta a etapa do funil com maior drop-off"),
    ("Hypothesis", "Hipótese testável (Se/Então/Porque)",
     "baseline de conversão vem do MTA · efeito mínimo detectável dimensionado a partir do lift do MMM"),
    ("Priorization", "ICE Score (Impact × Confidence × Ease)",
     "Impact ponderado pela contribuição incremental do MMM · Confidence sobe se já há geo-experiment"),
    ("Design", "Especificação do teste (Control/Variant, split, métrica)",
     "métrica primária = KPI do MMM · tamanho de amostra calculado na Calculadora A/B"),
    ("Tracking", "Pipeline operacional (Ideation → Running)",
     "eventos do Data Layer instrumentam a variante e conectam ao QA"),
    ("Results", "Resultado estatístico (lift, p-value, PBB)",
     "comparado contra o crédito da atribuição e contra a curva de resposta do MMM"),
    ("Learning", "Insight e próximo teste",
     "retroalimenta o prior do MMM, revalida o peso do MTA e gera novo item em Opportunity"),
]
